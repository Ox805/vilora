"""Verify mediation/file_extraction.extract() against tiny in-memory fixtures.

Run from the project root:
    ./venv/bin/python scripts/verify_file_extraction.py
"""
import io
import os
import sys
import base64

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from mediation.file_extraction import extract, ExtractionResult


def _docx_fixture():
    from docx import Document
    doc = Document()
    doc.add_paragraph("Hello from a docx fixture.")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _xlsx_fixture():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "TestSheet"
    ws['A1'] = "Header"
    ws['A2'] = "value"
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pptx_fixture():
    from pptx import Presentation
    p = Presentation()
    slide = p.slides.add_slide(p.slide_layouts[5])
    slide.shapes.title.text = "Hello from a pptx fixture."
    buf = io.BytesIO()
    p.save(buf)
    return buf.getvalue()


def _pdf_fixture():
    return base64.b64decode(
        b'JVBERi0xLjQKMSAwIG9iago8PC9UeXBlL0NhdGFsb2cvUGFnZXMgMiAwIFI+PgplbmRvYmoKMiAw'
        b'IG9iago8PC9UeXBlL1BhZ2VzL0NvdW50IDEvS2lkc1szIDAgUl0+PgplbmRvYmoKMyAwIG9iago8'
        b'PC9UeXBlL1BhZ2UvUGFyZW50IDIgMCBSL01lZGlhQm94WzAgMCA2MTIgNzkyXS9SZXNvdXJjZXMg'
        b'PDwvRm9udDw8L0YxIDQgMCBSPj4+Pi9Db250ZW50cyA1IDAgUj4+CmVuZG9iago0IDAgb2JqCjw8'
        b'L1R5cGUvRm9udC9TdWJ0eXBlL1R5cGUxL0Jhc2VGb250L0hlbHZldGljYT4+CmVuZG9iago1IDAg'
        b'b2JqCjw8L0xlbmd0aCA0ND4+c3RyZWFtCkJUIC9GMSAxOCBUZiAxMDAgNzAwIFRkIChIZWxsbyBQ'
        b'REYpIFRqIEVUCmVuZHN0cmVhbQplbmRvYmoKeHJlZgowIDYKMDAwMDAwMDAwMCA2NTUzNSBmIAow'
        b'MDAwMDAwMDA5IDAwMDAwIG4gCjAwMDAwMDAwNTYgMDAwMDAgbiAKMDAwMDAwMDExMSAwMDAwMCBu'
        b'IAowMDAwMDAwMjAyIDAwMDAwIG4gCjAwMDAwMDAyNjEgMDAwMDAgbiAKdHJhaWxlcgo8PC9TaXpl'
        b'IDYvUm9vdCAxIDAgUj4+CnN0YXJ0eHJlZgozNTYKJSVFT0YK'
    )


def _png_fixture():
    return base64.b64decode(
        b'iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII='
    )


def check(label, content_type, blob_bytes, expected_kind, must_contain=None):
    result = extract(content_type, blob_bytes)
    assert result.kind == expected_kind, f"{label}: expected {expected_kind}, got {result.kind} (error={result.error})"
    if expected_kind == 'text' and must_contain:
        assert must_contain in (result.text or ''), f"{label}: missing '{must_contain}' in extracted text"
    if expected_kind == 'image':
        assert result.image_b64, f"{label}: image_b64 empty"
        assert result.image_media_type == content_type, f"{label}: wrong media type"
    print(f"  OK  {label} -> kind={result.kind}, truncated={result.was_truncated}")


def main():
    print("Running file_extraction verifier...")

    check("plain text",   'text/plain',    b"hello world",      'text', must_contain='hello')
    check("markdown",     'text/markdown', b"# title\nbody",    'text', must_contain='title')
    check("csv",          'text/csv',      b"a,b,c\n1,2,3",     'text', must_contain='a,b,c')

    check("docx",         'application/vnd.openxmlformats-officedocument.wordprocessingml.document', _docx_fixture(), 'text', must_contain='Hello from a docx')
    check("xlsx",         'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',       _xlsx_fixture(), 'text', must_contain='TestSheet')
    check("pptx",         'application/vnd.openxmlformats-officedocument.presentationml.presentation', _pptx_fixture(), 'text', must_contain='Hello from a pptx')

    check("pdf",          'application/pdf', _pdf_fixture(), 'text', must_contain='Hello PDF')

    check("png",          'image/png',  _png_fixture(), 'image')
    check("jpeg-shape",   'image/jpeg', b'\xff\xd8\xff\xd9', 'image')

    check("zip-unsupported", 'application/zip',  b"PK\x03\x04...",        'unreadable')
    check("empty file",      'text/plain',       b"",                     'unreadable')
    check("non-utf8 text",   'text/plain',       b"\xff\xfe\xff\xfe",     'unreadable')

    large = b"a" * 50000
    result = extract('text/plain', large)
    assert result.kind == 'text', f"large-text: expected text, got {result.kind}"
    assert result.was_truncated, "large-text: expected was_truncated=True"
    assert len(result.text) == 40000, f"large-text: expected 40000 chars, got {len(result.text)}"
    print(f"  OK  large text truncates at 40000 chars")

    print("\nAll checks passed.")


if __name__ == '__main__':
    main()
