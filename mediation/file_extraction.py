"""Pure file-extraction module: bytes -> ExtractionResult.

Dispatches by content type. Never raises -- any failure surfaces as
ExtractionResult(kind='unreadable', error=...).
"""
import base64
import io
import sys
from dataclasses import dataclass
from typing import Optional

PER_FILE_CHAR_CAP = 40000


@dataclass
class ExtractionResult:
    kind: str
    text: Optional[str] = None
    image_b64: Optional[str] = None
    image_media_type: Optional[str] = None
    was_truncated: bool = False
    error: Optional[str] = None


def _truncate(text):
    if len(text) > PER_FILE_CHAR_CAP:
        return text[:PER_FILE_CHAR_CAP], True
    return text, False


def _extract_text(blob_bytes):
    try:
        return blob_bytes.decode('utf-8')
    except UnicodeDecodeError as e:
        raise ValueError(f"not valid UTF-8: {e}")


def _extract_pdf(blob_bytes):
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(blob_bytes))
    parts = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        if page_text:
            parts.append(page_text)
    return "\n\n".join(parts)


def _extract_docx(blob_bytes):
    from docx import Document
    doc = Document(io.BytesIO(blob_bytes))
    parts = []
    for para in doc.paragraphs:
        if para.text:
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _extract_xlsx(blob_bytes):
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(blob_bytes), data_only=True, read_only=True)
    try:
        parts = []
        for sheet in wb.worksheets:
            parts.append(f'Sheet "{sheet.title}":')
            for row in sheet.iter_rows(values_only=True):
                cells = [str(v) for v in row if v is not None]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    finally:
        wb.close()


def _extract_pptx(blob_bytes):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(blob_bytes))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"Slide {i}:")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


_TEXT_TYPES = {'text/plain', 'text/markdown', 'text/csv'}
_IMAGE_TYPES = {'image/jpeg', 'image/png', 'image/gif', 'image/webp'}

_DOC_EXTRACTORS = {
    'application/pdf': _extract_pdf,
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': _extract_docx,
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': _extract_xlsx,
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': _extract_pptx,
}


def extract(content_type, blob_bytes):
    """Turn bytes into an ExtractionResult based on content type. Never raises."""
    if not blob_bytes:
        return ExtractionResult(kind='unreadable', error='empty file')

    try:
        if content_type in _IMAGE_TYPES:
            return ExtractionResult(
                kind='image',
                image_b64=base64.b64encode(blob_bytes).decode('ascii'),
                image_media_type=content_type,
            )
        if content_type in _TEXT_TYPES:
            text = _extract_text(blob_bytes)
        elif content_type in _DOC_EXTRACTORS:
            text = _DOC_EXTRACTORS[content_type](blob_bytes)
        else:
            return ExtractionResult(kind='unreadable', error='unsupported content type')

        text = text.strip()
        if not text:
            return ExtractionResult(kind='unreadable', error='no text content (possibly a scanned document)')

        clipped, truncated = _truncate(text)
        return ExtractionResult(kind='text', text=clipped, was_truncated=truncated)

    except Exception as e:
        sys.stderr.write(f"[file_extraction] {content_type}: {e}\n")
        return ExtractionResult(kind='unreadable', error=str(e)[:200])
